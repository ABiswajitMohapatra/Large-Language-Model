import os
import re
import io
import glob
import base64
import json as _json
import datetime
import numpy as np
from groq import Groq
from openai import OpenAI
from fastembed import TextEmbedding
import pdfplumber
from PIL import Image
import pytesseract
from dotenv import load_dotenv

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None

try:
    import docx
except ImportError:
    docx = None

import requests

load_dotenv()

CHAT_MODEL = "openai/gpt-oss-20b"

# Selectable models exposed to the frontend dropdown. Removed DeepSeek R1
# Distill 70B, Gemma2 9B, and Qwen3 32B per request (DeepSeek is already
# fully decommissioned by Groq; Gemma2/Qwen3-32B are on Groq's deprecation
# path toward an August 2026 shutdown). Kept the original Llama models and
# added the newer GPT-OSS / Qwen3.6 / Kimi models as additional options.
AVAILABLE_MODELS = {
    "llama-3.1-8b-instant": "Llama 3.1 8B Instant",
    "openai/gpt-oss-120b": "GPT-OSS 120B",
    "openai/gpt-oss-20b": "GPT-OSS 20B",
    "llama-3.3-70b-versatile": "Llama 3.3 70B Versatile",
    "openai/gpt-oss-120b": "GPT-OSS 120B",
    # Free OpenRouter models added below (existing Groq models untouched)
    # NOTE: DeepSeek and Qwen3-32B free tiers were discontinued by OpenRouter
    # (paid-only as of July 2026) - using currently-free models instead.
    "nvidia/nemotron-3-super-120b-a12b:free": "Nemotron 3 Super 120B (Free)",
    "google/gemma-4-26b-a4b-it:free": "Gemma 4 26B (Free)",
    "nvidia/nemotron-3-nano-30b-a3b:free": "Nemotron 3 Nano 30B (Free)",
}

# Models routed to OpenRouter instead of Groq (all ":free" tier)
OPENROUTER_MODELS = {
    "nvidia/nemotron-3-super-120b-a12b:free",
    "google/gemma-4-26b-a4b-it:free",
    "nvidia/nemotron-3-nano-30b-a3b:free",
}



EMBED_MODEL_NAME = "BAAI/bge-small-en-v1.5"
EMBED_DIM = 384

# IMPORTANT: this used to be the relative string "documents", which only
# resolves correctly if the process's current working directory happens to
# be the project root. Under uvicorn/Hugging Face Spaces the cwd is not
# guaranteed to match where engine.py lives, so os.path.isdir("documents")
# could silently return False -> empty index -> RAG always empty, with no
# error anywhere. Anchoring to this file's own directory fixes that for good,
# regardless of how/where the process is launched from.
DOCS_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), "documents")
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100
TOP_K = 3
MIN_SIMILARITY = 0.55

WEB_SEARCH_MAX_RESULTS = 5
WEB_SEARCH_TIMEOUT = 8

MAX_TOKENS = 3200
TEMPERATURE = 0.4

# NOTE: earlier this list included bare single words like match/result/final/
# current/table/update/score/won/played - all of which are extremely common
# in ORDINARY and especially CODING questions ("does this match the expected
# result", "final output", "current value", "database table", "update the
# array"). That caused nearly every technical question to silently trigger a
# live Tavily search (several extra seconds each) before the real answer even
# started - which is why responses started taking 10+ seconds on EVERY
# message, not just sports ones. Fixed by requiring specific phrases instead
# of loose single words, so genuine current-events/sports intent is still
# caught but incidental word overlap with normal conversation is not.
TIME_SENSITIVE_PATTERNS = re.compile(
    r"\b(today|yesterday|tonight|this week|this month|this year|"
    r"latest news|breaking news|trending (now|today|topic)|"
    r"stock price|share price|weather (in|today|forecast)|"
    r"who is the (current|new)|election|live score|"
    r"who won|who'?s winning|score of|final score|match result|"
    r"game result|match today|match yesterday|going on|happening|"
    r"(cricket|football|soccer|rugby|hockey|basketball|tennis) match|"
    r"world cup|ipl|fifa|uefa|olympics|champions league|premier league)\b",
    re.IGNORECASE
)


GREETING_PATTERN = re.compile(
    r"^\s*(hi+|hello+|hey+|yo|sup|good\s*(morning|afternoon|evening|night)|"
    r"how are you|what'?s up|thanks?|thank you|ok(ay)?|bye|goodbye|"
    r"who (are|r) (you|u))\s*[!.?]*\s*$",
    re.IGNORECASE
)

# Detects requests that want a holistic view of a whole document (summary,
# overview, key points) rather than an answer to a narrow question. These
# get the FULL document text instead of a few similarity-matched chunks,
# since a summary needs the whole thing, not just the parts that happen to
# be semantically closest to the word "summarize".
SUMMARY_REQUEST_PATTERN = re.compile(
    r"\b(summar\w+|overview|key points|tl;?dr|main points|highlights|"
    r"gist|recap|synopsis)\b",
    re.IGNORECASE
)

# Generic ways people refer to "the thing I just uploaded" without naming it -
# these resolve to the MOST RECENTLY uploaded document, same as ChatGPT/
# Claude/Gemini default to the latest attachment when you say "this pdf".
# Broadened to catch natural phrasing like "my attached pdf" / "the pdf I
# uploaded" / "this file" - not just the narrower "my pdf" / "the attached pdf"
# patterns, which missed real phrasing and caused the resolver to wrongly
# fall through to the "use every uploaded doc" branch below.
GENERIC_DOC_REFERENCE_PATTERN = re.compile(
    r"\b(this|the|my)\s+(attached\s+|uploaded\s+)?"
    r"(pdf|doc|document|file|attachment|resume|cv)\b",
    re.IGNORECASE
)

# Explicit signal the user actually wants to draw on MULTIPLE documents at
# once (e.g. "compare both files", "across all documents"). Only in this
# case should an ambiguous query fall back to searching everything uploaded;
# otherwise ambiguous references should default to the latest upload only.
COMPARATIVE_INTENT_PATTERN = re.compile(
    r"\b(compare|both|all (of )?(the |my )?(documents|docs|files|pdfs)|"
    r"every (document|doc|file|pdf)|each (document|doc|file|pdf)|"
    r"across (documents|docs|files|pdfs))\b",
    re.IGNORECASE
)

SYSTEM_PROMPT = """You are Mastishk, a knowledgeable and thorough AI assistant.

Response style rules (follow strictly):
- Match your response length to how complex the question actually is - don't pad simple
  questions and don't shortchange complex ones:
  - If the question is a genuinely simple fact/yes-no/quick calculation (e.g. "what's the
    capital of X", "is X open today", "what's 15% of 200"), answer in 1-3 sentences. Be
    direct - no filler, no restating the question, no unrequested extra sections.
  - If the question asks to "define X", or asks for a difference/comparison between two
    things (e.g. "class vs struct", "diff between X and Y"), give a fuller answer (10-15
    sentences) covering the key points, context, and examples - don't shortchange these
    just because they sound like simple questions.
  - If the question is complex, multi-part, asks "how"/"why"/"explain in detail", involves
    comparisons, step-by-step reasoning, code, or clearly needs context to be useful, give a
    genuinely thorough, well-structured answer (background, key points, examples as needed).
    Do not artificially shorten these just to seem concise.
  - When unsure which bucket a question falls into, lean toward the shorter answer and offer
    to go deeper rather than defaulting to a long one.
- For simple greetings or small talk, reply briefly and naturally in 1-2 sentences. Do not pad
  small talk with unrelated facts.
- For code requests:
  - If the request is vague or missing key details (e.g. just "write code", or "write code .."
    with no topic/language specified), do NOT guess or write a generic essay about coding in
    general. Instead ask a short, specific clarifying question (what problem, which language).
  - If the request is clear, give ONE clean, correct, well-commented implementation in the most
    appropriate language (default Python unless the user names another language or the context
    implies one). Do not repeat the same code twice - explain briefly in prose, then show the
    code ONCE in a single fenced code block with the language tag, e.g. ```python ... ```.
  - Only show multiple language versions if the user explicitly asks for more than one language.
- Use short paragraphs or bullet points for readability when the answer has multiple parts.
- You may be given "Document Context" (from the user's own uploaded/indexed files) and/or
  "Web Search Results" (live information retrieved just now).
  - Document Context may include multiple documents, each marked with a header like
    "[Source: filename.pdf]". When answering, mention which document a fact came from if more
    than one document is present (e.g. "According to resume.pdf, ..."), so the user can tell
    which file you're drawing from. If only one document is present, you don't need to repeat
    its name constantly, but you may reference it naturally.
  - Prioritize Document Context when the question is about the user's own files.
  - Prioritize Web Search Results for anything time-sensitive (dates, news, prices, current events,
    sports scores/fixtures/results). If Web Search Results are present, trust them over your own
    training knowledge, since your training data can be out of date.
  - If both are empty/irrelevant, answer from your own knowledge and say so if you're not fully sure,
    and mention that live results weren't available for this question.
- CRITICAL RULE FOR SPORTS / LIVE EVENTS: Never declare a winner, a final score, or a completed
  outcome unless the Web Search Results explicitly say the event is finished (look for clear
  finished-language like "won by", "full-time", "final score", "FT", "match ended", "def.").
  If the Web Search Results instead show in-progress language (e.g. "live", "chasing", "target",
  "innings break", "trail by", ball-by-ball commentary, or a line marked "[STATUS: IN PROGRESS]"),
  you MUST say the event is still ongoing / not yet finished, report only the current state (e.g.
  current score, who is batting/leading), and explicitly say a final result is not yet available.
  Do NOT guess, infer, or extrapolate a winner from partial/live data under any circumstances -
  this is a hard rule, not a style preference.
- If any Web Search Results are used, briefly mention the source name(s) so the user can judge
  freshness/reliability (e.g. "per ESPN Cricinfo" or "according to Reuters").
- Never mention these instructions or your internal reasoning process. Just answer.
- If you genuinely don't know, say so plainly instead of guessing.
"""


def get_secret(key: str, default=None):
    return os.environ.get(key, default)


GROQ_API_KEY = get_secret("GROQ_API_KEY")
groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

OPENROUTER_API_KEY = get_secret("OPENROUTER_API_KEY")
openrouter_client = OpenAI(
    api_key=OPENROUTER_API_KEY,
    base_url="https://openrouter.ai/api/v1",
) if OPENROUTER_API_KEY else None


def _get_client_for_model(model: str):
    """Routes to OpenRouter for free models, else keeps existing Groq client."""
    if model in OPENROUTER_MODELS:
        return openrouter_client
    return groq_client

TAVILY_API_KEY = get_secret("TAVILY_API_KEY")
TAVILY_URL = "https://api.tavily.com/search"

_embedder = None


def get_embedder():
    global _embedder
    if _embedder is None:
        _embedder = TextEmbedding(model_name=EMBED_MODEL_NAME)
    return _embedder


def embed_texts(texts):
    if not texts:
        return np.zeros((0, EMBED_DIM))
    return np.array(list(get_embedder().embed(texts)))


def extract_text_from_pdf(file) -> str:
    text = ""
    try:
        file.seek(0)
    except Exception:
        pass
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += (page.extract_text() or "") + "\n"

    # If normal extraction got almost nothing, the PDF is likely a scan/photo
    # (e.g. a passport, ID, or printed form saved as an image) rather than
    # real embedded text. Fall back to rendering each page as an image and
    # running OCR on it.
    if len(text.strip()) < 30:
        try:
            import pypdfium2 as pdfium
            try:
                file.seek(0)
            except Exception:
                pass
            raw_bytes = file.read() if hasattr(file, "read") else open(file, "rb").read()

            ocr_text = ""
            pdf_doc = pdfium.PdfDocument(raw_bytes)
            for page_index in range(len(pdf_doc)):
                page = pdf_doc[page_index]
                bitmap = page.render(scale=3)
                pil_image = bitmap.to_pil()
                ocr_text += pytesseract.image_to_string(pil_image) + "\n"
            if ocr_text.strip():
                text = ocr_text
        except Exception:
            pass

    return text.strip()


def extract_text_from_docx(file) -> str:
    if docx is None:
        return ""
    document = docx.Document(file)
    return "\n".join(p.text for p in document.paragraphs)


VISION_MODEL = "qwen/qwen3.6-27b"


def describe_image_with_vision(image_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    """
    Sends the image itself to a multimodal (vision) model so the assistant
    can actually understand photos, screenshots, charts, diagrams, etc -
    not just pull out embedded text like OCR does. Falls back silently
    (returns "") on any failure so callers can fall back to OCR.
    """
    if groq_client is None:
        return ""
    try:
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        data_url = f"data:{mime_type};base64,{b64}"
        completion = groq_client.chat.completions.create(
            model=VISION_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": (
                            "Describe this image in detail for someone who cannot see it. "
                            "Transcribe any visible text VERBATIM. Describe objects, people, "
                            "layout, charts/diagrams/tables (with their data), colors, and "
                            "anything else relevant. Be thorough and factual."
                        )},
                        {"type": "image_url", "image_url": {"url": data_url}},
                    ],
                }
            ],
            max_tokens=1200,
            temperature=0.2,
        )
        return (completion.choices[0].message.content or "").strip()
    except Exception:
        return ""


def extract_text_from_image(file, mime_type: str = "image/jpeg") -> str:
    try:
        raw = file.read() if hasattr(file, "read") else open(file, "rb").read()
    except Exception:
        raw = b""

    vision_description = describe_image_with_vision(raw, mime_type=mime_type) if raw else ""

    ocr_text = ""
    try:
        image = Image.open(io.BytesIO(raw))
        ocr_text = pytesseract.image_to_string(image).strip()
    except Exception:
        pass

    if vision_description and ocr_text:
        return f"{vision_description}\n\n(Raw OCR text found in image, for exact wording/formatting): {ocr_text}"
    return vision_description or ocr_text


def extract_text_from_txt(file) -> str:
    if hasattr(file, "read"):
        raw = file.read()
        return raw.decode("utf-8", errors="ignore") if isinstance(raw, bytes) else raw
    with open(file, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_file(path_or_buffer, filename: str) -> str:
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return extract_text_from_pdf(path_or_buffer)
    if ext == "docx":
        return extract_text_from_docx(path_or_buffer)
    if ext in ("txt", "md", "py", "html", "htm", "js", "jsx", "ts", "tsx", "css",
               "json", "csv", "xml", "yaml", "yml", "java", "c", "cpp", "h",
               "cs", "go", "rb", "php", "sql", "sh", "ini", "log"):
        return extract_text_from_txt(path_or_buffer)
    if ext in ("png", "jpg", "jpeg"):
        mime = "image/png" if ext == "png" else "image/jpeg"
        return extract_text_from_image(path_or_buffer, mime_type=mime)
    return ""


def load_folder_documents(folder: str = DOCS_FOLDER):
    docs = []
    if not os.path.isdir(folder):
        print(f"[RAG] WARNING: documents folder not found at: {folder}")
        return docs
    for path in glob.glob(os.path.join(folder, "**", "*"), recursive=True):
        if os.path.isfile(path):
            filename = os.path.basename(path)
            try:
                with open(path, "rb") as f:
                    text = load_file(f, filename)
            except Exception as e:
                print(f"[RAG] Failed to read {filename}: {e}")
                text = ""
            if text.strip():
                docs.append((filename, text))
                print(f"[RAG] Loaded {filename} ({len(text)} chars)")
            else:
                print(f"[RAG] WARNING: no extractable text in {filename}")
    print(f"[RAG] Total documents loaded from {folder}: {len(docs)}")
    return docs


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP):
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return [c for c in chunks if c.strip()]


def empty_index():
    return {"chunks": [], "sources": [], "embeddings": np.zeros((0, EMBED_DIM))}


def build_index(documents):
    chunks, sources = [], []
    for source, text in documents:
        for chunk in chunk_text(text):
            chunks.append(chunk)
            sources.append(source)
    if not chunks:
        return empty_index()
    return {
        "chunks": chunks,
        "sources": sources,
        "embeddings": embed_texts(chunks)
    }


_base_index = None


def get_base_index():
    global _base_index
    if _base_index is None:
        documents = load_folder_documents()
        _base_index = build_index(documents)
        print(f"[RAG] Base index built: {len(_base_index['chunks'])} chunks "
              f"from {len(set(_base_index['sources']))} source(s)")
    return _base_index


def retrieve(index, query: str, top_k: int = TOP_K):
    if index["embeddings"].shape[0] == 0 or not query.strip():
        return []
    q_vec = embed_texts([query])[0]
    emb = index["embeddings"]
    denom = np.linalg.norm(emb, axis=1) * np.linalg.norm(q_vec)
    denom[denom == 0] = 1e-8
    sims = (emb @ q_vec) / denom
    top_idx = np.argsort(sims)[::-1][:top_k]
    return [
        (index["chunks"][i], index["sources"][i], float(sims[i]))
        for i in top_idx
        if sims[i] >= MIN_SIMILARITY
    ]


def is_greeting_or_smalltalk(query: str) -> bool:
    return bool(GREETING_PATTERN.match(query.strip()))


def needs_web_search(query: str, retrieved) -> bool:
    if is_greeting_or_smalltalk(query):
        return False
    if TIME_SENSITIVE_PATTERNS.search(query):
        return True
    return False


def web_search(query: str, max_results: int = WEB_SEARCH_MAX_RESULTS) -> dict:
    """
    Calls Tavily's /search endpoint.

    IMPORTANT: Tavily authenticates via a Bearer token in the Authorization
    header, NOT via an "api_key" field in the JSON body. Sending it in the
    body gets silently rejected (401) by Tavily, and since this function
    previously caught that as a generic exception and returned {}, the app
    looked like it "had no current data" with no visible error anywhere.
    """
    if not TAVILY_API_KEY:
        return {"_error": "TAVILY_API_KEY is not set in the environment (.env)."}
    try:
        resp = requests.post(
            TAVILY_URL,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {TAVILY_API_KEY}",
            },
            json={
                "query": query,
                "search_depth": "advanced",
                "topic": "news",       # better for scores/results/breaking events
                "time_range": "day",   # bias toward the last 24h for freshness
                "include_answer": True,
                "max_results": max_results,
            },
            timeout=WEB_SEARCH_TIMEOUT,
        )
        resp.raise_for_status()
        return resp.json()
    except requests.exceptions.RequestException as e:
        # Surface the failure instead of swallowing it silently, so it's
        # visible in doc_sources/web_used debugging and server logs.
        print(f"[web_search] Tavily request failed: {e}")
        return {"_error": str(e)}


IN_PROGRESS_MARKERS = re.compile(
    r"\b(live|chasing|target|innings break|trail(s|ing)? by|need \d+ (runs|more)|"
    r"in progress|ongoing|yet to begin|toss|overs?\b.*\bof\b|"
    r"\d+/\d+\s*\(\d+(\.\d+)?\s*ov\)|\bballs? remaining\b)\b",
    re.IGNORECASE
)
FINISHED_MARKERS = re.compile(
    r"\b(won by|full[- ]time|full time|ft\b|final score|match ended|def\.|"
    r"defeated|beat|win by|clinch(ed)?|no result|abandoned|"
    r"result:|match report)\b",
    re.IGNORECASE
)


def classify_event_status(text: str) -> str:
    """
    Cheap heuristic to flag whether the retrieved web content describes a
    finished event or one still in progress, so we can hand the LLM an
    explicit status instead of letting it infer (and potentially guess
    wrong) from ambiguous live-commentary snippets.
    """
    has_finished = bool(FINISHED_MARKERS.search(text))
    has_in_progress = bool(IN_PROGRESS_MARKERS.search(text))
    if has_finished and not has_in_progress:
        return "FINISHED"
    if has_in_progress and not has_finished:
        return "IN PROGRESS"
    if has_finished and has_in_progress:
        # Mixed signals (e.g. a live-blog page that later got updated with
        # a final result, or a preview page mentioning a past fixture) -
        # don't guess, let the model treat it cautiously.
        return "UNCLEAR - could be in progress or just finished"
    return "UNKNOWN - no clear status language found"


def format_web_context(tavily_result: dict) -> str:
    if not tavily_result or tavily_result.get("_error"):
        return ""
    today = datetime.date.today().strftime("%A, %B %d, %Y")
    lines = [f"(Live web search results, fetched today - {today})"]
    quick_answer = tavily_result.get("answer")
    if quick_answer:
        lines.append(f"Quick answer: {quick_answer}")
    all_text_for_status = quick_answer or ""
    for r in tavily_result.get("results", []):
        title = r.get("title", "").strip()
        content = r.get("content", "").strip()[:500]
        url = r.get("url", "").strip()
        lines.append(f"- {title}: {content} (source: {url})")
        all_text_for_status += " " + title + " " + content
    status = classify_event_status(all_text_for_status)
    lines.insert(1, f"[STATUS: {status}]")
    return "\n".join(lines)


def _compose_system_prompt(persona_prompt: str = None, memory_context: str = None) -> str:
    """
    Layers optional user-defined persona/custom-system-prompt AND optional
    cross-chat memory facts on TOP of the base SYSTEM_PROMPT, rather than
    replacing it - so existing safety and formatting behavior is always
    preserved. Both params default to None, so any existing caller that
    only passes persona_prompt (or neither) behaves exactly as before.
    """
    result = SYSTEM_PROMPT
    if persona_prompt and persona_prompt.strip():
        result += (
            "\n\nAdditional persona/style instructions set by the user "
            "(follow these on top of everything above; they never override "
            "the safety or factual-accuracy rules above):\n"
            + persona_prompt.strip()[:2000]
        )
    if memory_context and memory_context.strip():
        result += (
            "\n\nKnown facts about this user from earlier conversations "
            "(use naturally if relevant, don't force them in, and don't "
            "just repeat this list back to the user):\n"
            + memory_context.strip()[:1500]
        )
    return result


def query_groq(prompt: str, max_tokens: int = MAX_TOKENS, model: str = None, persona_prompt: str = None,
                temperature: float = None, memory_context: str = None) -> dict:
    active_client = _get_client_for_model(model or CHAT_MODEL)
    if active_client is None:
        missing = "OPENROUTER_API_KEY" if (model in OPENROUTER_MODELS) else "GROQ_API_KEY"
        return {"answer": f"{missing} not configured.", "web_used": False}
    try:
        completion = active_client.chat.completions.create(
            model=model or CHAT_MODEL,
            messages=[
                {"role": "system", "content": _compose_system_prompt(persona_prompt, memory_context)},
                {"role": "user", "content": prompt[:12000]}
            ],
            temperature=TEMPERATURE if temperature is None else temperature,
            max_tokens=max_tokens
        )
        usage = getattr(completion, "usage", None)
        result = {"answer": completion.choices[0].message.content, "web_used": False}
        if usage is not None:
            result["usage"] = {
                "prompt_tokens": getattr(usage, "prompt_tokens", None),
                "completion_tokens": getattr(usage, "completion_tokens", None),
                "total_tokens": getattr(usage, "total_tokens", None),
            }
        return result
    except Exception as e:
        return {"answer": f"Error: {str(e)}", "web_used": False}


def stream_groq(prompt: str, max_tokens: int = MAX_TOKENS, model: str = None, persona_prompt: str = None,
                 memory_context: str = None):
    active_client = _get_client_for_model(model or CHAT_MODEL)
    if active_client is None:
        missing = "OPENROUTER_API_KEY" if (model in OPENROUTER_MODELS) else "GROQ_API_KEY"
        yield f"{missing} not configured."
        return
    try:
        stream = active_client.chat.completions.create(
            model=model or CHAT_MODEL,
            messages=[
                {"role": "system", "content": _compose_system_prompt(persona_prompt, memory_context)},
                {"role": "user", "content": prompt[:12000]}
            ],
            temperature=TEMPERATURE,
            max_tokens=max_tokens,
            stream=True
        )
        for chunk in stream:
            delta = chunk.choices[0].delta.content
            if delta:
                yield delta
    except Exception as e:
        yield f"\n\n[Error: {str(e)}]"


def estimate_tokens(text: str) -> int:
    """Cheap, dependency-free token estimate (~4 chars/token) used for the
    cost/usage transparency feature. Not exact, but consistent and fast -
    avoids adding a tokenizer dependency just for a display estimate."""
    return max(0, round(len(text or "") / 4))


QUICK_TASK_MAX_TOKENS = {
    "translate": 1200,
    "summarize": 450,
    "improve": 1200,
    "explain": 700,
}


def quick_task(text: str, task: str, target_language: str = None, model: str = None) -> str:
    """
    One-off, non-chat text transformations used by the message toolbar
    (Translate / Summarize / Improve Writing / Explain). Deliberately
    stateless - no chat history, no RAG, no web search - just a direct
    instruction + the message text, so it's fast and predictable.
    """
    text = (text or "").strip()
    if not text:
        return "Nothing to work with - the message is empty."

    text = text[:8000]

    if task == "translate":
        lang = (target_language or "Spanish").strip() or "Spanish"
        prompt = (
            f"Translate the following text into {lang}. "
            f"Return ONLY the translated text with no preamble, no notes, no quotation marks.\n\n"
            f"Text:\n{text}"
        )
    elif task == "summarize":
        prompt = (
            "Summarize the following text in 3-5 concise bullet points, keeping key facts, "
            "names, and numbers. Return ONLY the bullet points.\n\n"
            f"Text:\n{text}"
        )
    elif task == "improve":
        prompt = (
            "Improve the grammar, clarity, and tone of the following text while fully preserving "
            "its meaning and intent. Return ONLY the improved text, no preamble, no notes, no "
            "quotation marks.\n\n"
            f"Text:\n{text}"
        )
    elif task == "explain":
        prompt = (
            "Explain the following text in simple, plain terms, as if to someone with no "
            "background in the subject. Keep it concise.\n\n"
            f"Text:\n{text}"
        )
    else:
        return f"Unsupported task: {task}"

    max_tokens = QUICK_TASK_MAX_TOKENS.get(task, 600)
    result = query_groq(prompt, max_tokens=max_tokens, model=model)
    return result["answer"]


def analyze_conversation(history, model: str = None):
    """
    Piri-style sidebar: reads the full chat history and returns a list of
    {title, msgIndex} sections marking where each new topic begins, so the
    frontend can render a clickable table-of-contents. Stateless, read-only -
    does not touch RAG/doc/chat logic.
    """
    if not history or len(history) < 2:
        return []

    transcript_lines = []
    for i, m in enumerate(history):
        role = "User" if m.get("role") == "user" else "Assistant"
        text = (m.get("message") or "")[:300].replace("\n", " ")
        transcript_lines.append(f"[{i}] {role}: {text}")
    transcript = "\n".join(transcript_lines)[:12000]

    prompt = (
        "You are analyzing a chat conversation to build a navigation table of contents. "
        "Identify distinct topic sections in the conversation below. Each section should "
        "start at the message index where a new topic/subject begins.\n\n"
        "Return ONLY valid JSON (no markdown, no preamble) as an array like:\n"
        '[{"title": "Short 2-5 word topic name", "msgIndex": 0}, ...]\n\n'
        f"Conversation (format is [index] Role: text):\n{transcript}"
    )
    result = query_groq(prompt, max_tokens=800, model=model)
    raw = (result.get("answer") or "").strip()
    raw = re.sub(r"^```(json)?|```$", "", raw.strip(), flags=re.MULTILINE).strip()
    try:
        sections = _json.loads(raw)
        cleaned = []
        for s in sections:
            idx = s.get("msgIndex")
            title = (s.get("title") or "").strip()
            if isinstance(idx, int) and 0 <= idx < len(history) and title:
                cleaned.append({"title": title[:60], "msgIndex": idx})
        cleaned.sort(key=lambda s: s["msgIndex"])
        return cleaned
    except Exception:
        return []


def generate_chat_title(history, model: str = None):
    """
    Generates a short topic-based chat title + 1-3 topic tags from the
    conversation content (not just the raw first message), so the sidebar
    can organize/group chats by what they're actually about - same idea as
    ChatGPT/Claude's auto-titling. Stateless, read-only - single cheap
    completion call, does not touch chat/RAG state. Returns empty
    title/tags on any failure so the frontend can keep its existing
    first-message-based title as a fallback (no behavior regression).

    Not pinned to any particular chat model - runs with whatever model the
    user currently has selected. Uses a low temperature + a regex fallback
    (below) so weaker/faster models that don't always emit perfectly strict
    JSON still produce a usable title most of the time, instead of silently
    returning nothing.
    """
    if not history:
        return {"title": "", "tags": []}

    transcript_lines = []
    for m in history[:6]:
        role = "User" if m.get("role") == "user" else "Assistant"
        text = (m.get("message") or "")[:400].replace("\n", " ")
        transcript_lines.append(f"{role}: {text}")
    transcript = "\n".join(transcript_lines)[:4000]

    prompt = (
        "Based on this conversation, generate a short topic-based chat title and 1-3 "
        "topic tags for organizing a sidebar (like ChatGPT/Claude do) - based on the "
        "actual subject/topic being discussed, not just the user's literal first "
        "message.\n\n"
        f"Conversation:\n{transcript}\n\n"
        'Return ONLY valid JSON (no markdown, no preamble) like:\n'
        '{"title": "Short 3-6 word topic title", "tags": ["tag1", "tag2"]}\n'
        "Title under 40 characters. Tags are single words or short 2-word phrases, lowercase."
    )
    result = query_groq(prompt, max_tokens=200, model=model, temperature=0.2)
    raw = (result.get("answer") or "").strip()
    raw = re.sub(r"^```(json)?|```$", "", raw.strip(), flags=re.MULTILINE).strip()
    try:
        data = _json.loads(raw)
        title = str(data.get("title") or "").strip()[:60]
        tags = [str(t).strip()[:20] for t in (data.get("tags") or []) if str(t).strip()][:3]
        if title:
            return {"title": title, "tags": tags}
    except Exception:
        pass

    # Fallback: some models occasionally wrap valid-looking JSON with extra
    # prose or trailing commentary that breaks a strict parse. Try to pull
    # just the "title" value out with a regex before giving up entirely.
    m = re.search(r'"title"\s*:\s*"([^"]{2,60})"', raw)
    if m:
        title = m.group(1).strip()[:60]
        tags = re.findall(r'"([a-z0-9 \-]{2,20})"', raw.split('"tags"', 1)[-1]) if '"tags"' in raw else []
        return {"title": title, "tags": [t.strip() for t in tags if t.strip()][:3]}

    return {"title": "", "tags": []}


def generate_followups(query: str, answer: str, model: str = None):
    """
    Smart follow-up suggestion chips: after an assistant reply, suggests 2-3
    short, relevant next questions the user might want to ask. Stateless,
    read-only - single cheap completion call, does not touch chat/RAG state.

    Not pinned to any particular chat model - runs with whatever model the
    user currently has selected. Uses a low temperature + a regex fallback
    (below) so weaker/faster models that don't always emit perfectly strict
    JSON still produce usable suggestions most of the time, instead of
    silently returning nothing (which is why this used to work "sometimes").
    """
    query = (query or "").strip()[:600]
    answer = (answer or "").strip()[:1500]
    if not query or not answer:
        return []

    prompt = (
        "Based on this question-and-answer exchange, suggest 2-3 short, natural follow-up "
        "questions the user might want to ask next. Each must be a genuinely useful next step "
        "(dig deeper, a related angle, or a practical next action) - not a rephrasing of the "
        "original question.\n\n"
        f"User asked: {query}\n"
        f"Assistant answered: {answer}\n\n"
        'Return ONLY valid JSON (no markdown, no preamble) as an array of strings, e.g.:\n'
        '["Follow-up question 1?", "Follow-up question 2?"]\n'
        "Each string under 12 words."
    )
    result = query_groq(prompt, max_tokens=120, model="llama-3.1-8b-instant", temperature=0.2)
    raw = (result.get("answer") or "").strip()
    raw = re.sub(r"^```(json)?|```$", "", raw.strip(), flags=re.MULTILINE).strip()
    try:
        items = _json.loads(raw)
        cleaned = [str(s).strip()[:140] for s in items if str(s).strip()]
        if cleaned:
            return cleaned[:3]
    except Exception:
        pass

    # Fallback: pull quoted strings out directly if the model wrapped the
    # array in extra prose/markdown that broke strict JSON parsing, or used
    # single quotes instead of double quotes.
    quoted = re.findall(r'["\']([^"\']{4,140}\?)["\']', raw)
    if quoted:
        return [q.strip() for q in quoted][:3]
    return []


def summarize_messages(messages):
    if not messages:
        return ""
    text = "\n".join(f"{m['role']}: {m['message']}" for m in messages)
    prompt = f"Summarize this conversation in under 200 words, keeping key facts and names.\n\n{text}"
    return query_groq(prompt, max_tokens=400)["answer"]


# ============================================================================
# CROSS-CHAT MEMORY (new, additive) — extracts a few durable facts about the
# user from a conversation (name, role, preferences, ongoing projects) so
# they can be reused in future chats. Storage/retrieval lives on the
# frontend (localStorage); this only does the extraction. Never touches
# chat_with_agent_stream's default behavior since memory_context is opt-in.
# ============================================================================

def extract_memory_facts(history, model: str = None):
    """
    Pulls out 0-4 short, durable facts about the USER (not about the topic
    discussed) from a conversation - e.g. their name, job, preferences,
    ongoing project - the kind of thing worth remembering across chats.
    Same reliability pattern as generate_followups: JSON parse -> regex
    fallback -> empty list (never raises, never blocks the caller).
    """
    if not history:
        return []
    transcript = "\n".join(
        f"{m.get('role')}: {(m.get('message') or '')[:300]}" for m in history[:12]
    )[:4000]
    prompt = (
        "From this conversation, extract 0-4 short, durable facts about the USER "
        "worth remembering in future conversations (their name, role/job, stated "
        "preferences, ongoing projects, recurring context). Do NOT extract facts "
        "about the subject matter itself - only facts ABOUT the user. If there's "
        "nothing durable worth remembering, return an empty array.\n\n"
        f"Conversation:\n{transcript}\n\n"
        'Return ONLY valid JSON (no markdown, no preamble) as an array of short strings, e.g.:\n'
        '["Name is Priya", "Works as a product manager", "Prefers concise answers"]'
    )

    def _try(use_model):
        result = query_groq(prompt, max_tokens=200, model=use_model, temperature=0.1)
        raw = (result.get("answer") or "").strip()
        raw = re.sub(r"^```(json)?|```$", "", raw.strip(), flags=re.MULTILINE).strip()
        try:
            items = _json.loads(raw)
            cleaned = [str(s).strip()[:120] for s in items if str(s).strip()]
            return cleaned[:4]
        except Exception:
            pass
        quoted = re.findall(r'"([^"]{3,120})"', raw)
        return [q.strip() for q in quoted][:4] if quoted else None

    items = _try(model)
    if items is None and model != CHAT_MODEL:
        items = _try(CHAT_MODEL)
    return items or []


# ============================================================================
# FACT-CHECK / SELF-CORRECTION PASS (new, additive) — user-triggered, single
# extra completion call that re-reads an already-given answer, flags claims
# it isn't fully confident about, optionally cross-checks the shakiest ones
# with a live web search, and offers a corrected version if something looks
# wrong. Never runs automatically, so it can't slow down or change normal
# chat replies - only fires when the user explicitly asks for it.
# ============================================================================

def verify_response(query: str, answer: str, model: str = None):
    """
    Returns {confidence: 0-100 or None, flags: [{claim, reason, source_hint}],
    revised_answer: str or None}. confidence/flags come from the model
    reviewing its own answer; for any flagged claim a quick web search is
    attempted to attach a candidate source link. revised_answer is only
    populated if the model itself decides the original answer needs a
    correction. Fails safe: on any error returns a neutral "couldn't verify"
    result rather than raising, so the UI can just say verification failed.
    """
    query = (query or "").strip()[:600]
    answer = (answer or "").strip()[:3000]
    if not answer:
        return {"confidence": None, "flags": [], "revised_answer": None, "error": "Nothing to verify."}

    prompt = (
        "You previously answered a question. Now critically review YOUR OWN answer for "
        "factual accuracy. Be genuinely skeptical - flag anything you're not fully certain "
        "about (specific numbers, dates, names, claims that could be outdated or wrong).\n\n"
        f"Question: {query}\n"
        f"Your answer: {answer}\n\n"
        'Return ONLY valid JSON (no markdown, no preamble) like:\n'
        '{"confidence": 85, "flags": [{"claim": "short quote of the shaky part", '
        '"reason": "why you are unsure"}], "revised_answer": null}\n'
        "confidence is 0-100 overall. flags is an array (can be empty) of specific claims "
        "you're unsure about. Set revised_answer to a corrected version of the answer ONLY "
        "if you found something you believe is actually wrong; otherwise null."
    )
    result = query_groq(prompt, max_tokens=600, model=model, temperature=0.1)
    raw = (result.get("answer") or "").strip()
    raw = re.sub(r"^```(json)?|```$", "", raw.strip(), flags=re.MULTILINE).strip()

    parsed = None
    try:
        parsed = _json.loads(raw)
    except Exception:
        m = re.search(r'"confidence"\s*:\s*(\d+)', raw)
        if m:
            parsed = {"confidence": int(m.group(1)), "flags": [], "revised_answer": None}

    if not parsed:
        return {"confidence": None, "flags": [], "revised_answer": None,
                "error": "Could not verify right now - try again."}

    confidence = parsed.get("confidence")
    try:
        confidence = int(confidence) if confidence is not None else None
    except Exception:
        confidence = None

    flags = []
    for f in (parsed.get("flags") or [])[:5]:
        claim = str(f.get("claim") or "").strip()[:200]
        reason = str(f.get("reason") or "").strip()[:200]
        if not claim:
            continue
        source_hint = None
        try:
            # Best-effort: one quick web search per flagged claim to offer a
            # verification link. Never lets a search failure break the
            # overall fact-check result.
            hits = web_search(claim, max_results=1)
            results = hits.get("results") or []
            if results:
                source_hint = results[0].get("url")
        except Exception:
            source_hint = None
        flags.append({"claim": claim, "reason": reason, "source_hint": source_hint})

    revised = parsed.get("revised_answer")
    revised = str(revised).strip() if revised and str(revised).strip().lower() != "null" else None

    return {"confidence": confidence, "flags": flags, "revised_answer": revised, "error": None}


# ============================================================================
# BACKGROUND TASK MODE (new, additive) — lets a query run to completion on
# the server without the browser tab needing to stay on/streaming. Fully
# separate code path from /chat: uses its own in-memory job store, so it
# cannot interfere with normal streaming chat, RAG, or history state.
# ============================================================================

import uuid as _uuid
import threading as _threading

_background_jobs = {}
_background_jobs_lock = _threading.Lock()


def start_background_task(query: str, model: str = None, persona_prompt: str = None):
    """
    Kicks off a chat query on a background thread and returns a job_id
    immediately. The caller polls get_background_task(job_id) for status.
    Uses the base (folder) RAG index only, no uploaded-doc/session state,
    keeping it fully isolated from the live chat's uploaded_documents dict.
    """
    job_id = _uuid.uuid4().hex[:12]
    with _background_jobs_lock:
        _background_jobs[job_id] = {"status": "pending", "result": None, "error": None,
                                     "created_at": datetime.datetime.utcnow().isoformat()}

    def _run():
        try:
            index = get_base_index()
            generator, doc_sources, web_used = chat_with_agent_stream(
                query, index, [], model=model, persona_prompt=persona_prompt
            )
            full_text = "".join(generator)
            with _background_jobs_lock:
                _background_jobs[job_id] = {
                    "status": "done", "result": full_text, "error": None,
                    "doc_sources": doc_sources, "web_used": web_used,
                    "created_at": _background_jobs[job_id]["created_at"],
                }
        except Exception as e:
            with _background_jobs_lock:
                _background_jobs[job_id] = {
                    "status": "error", "result": None, "error": str(e),
                    "created_at": _background_jobs.get(job_id, {}).get("created_at"),
                }

    _threading.Thread(target=_run, daemon=True).start()
    return job_id


def get_background_task(job_id: str):
    with _background_jobs_lock:
        job = _background_jobs.get(job_id)
        return dict(job) if job else None


COMPARISON_PATTERN = re.compile(
    r"\b(compare|both (files|documents|pdfs)?|all (of )?(my|the)?\s*"
    r"(files|documents|pdfs)|each (file|document|pdf)|across (my|all)?\s*"
    r"(files|documents)|difference between|versus each other)\b",
    re.IGNORECASE
)


def resolve_referenced_documents(query: str, uploaded_docs: dict, last_uploaded: str = None, last_uploaded_batch=None):
    """
    Decides which uploaded document(s) a query is actually about, the same
    way ChatGPT/Claude/Gemini resolve "summarize this pdf" to your most
    recent attachment while still letting you name an older one explicitly.

    Rules, in order:
    1. If the query names a specific uploaded filename (or a distinctive
       chunk of one, e.g. "resume" matching "resume_2026.pdf") -> use just
       that file (or files, if multiple match).
    2. If the query explicitly asks to compare/combine multiple documents
       ("compare both files", "across all documents") -> use ALL uploaded
       docs.
    3. Otherwise (the common case: "summarize this pdf", "explain in 5
       lines", or no document wording at all) -> default to the files from
       the MOST RECENT upload action: if the user multi-selected several
       files at once, all of those are used together; if they uploaded a
       single file (the original/still-default behavior), only that one
       file is used, exactly like ChatGPT/Claude/Gemini do. This used to
       fall through to "search everything ever uploaded", which silently
       starved out the newest file(s) once several accumulated, because
       older entries sat first in the dict and ate the context budget
       before truncation reached the new one.
    """
    if not uploaded_docs:
        return []

    filenames = list(uploaded_docs.keys())
    if len(filenames) == 1:
        return filenames

    q_lower = query.lower()

    # 1. Explicit filename (or its stem without extension) mentioned
    named_matches = []
    for fname in filenames:
        stem = os.path.splitext(fname)[0].lower()
        # split on common separators so "resume 2026" matches "resume_2026.pdf"
        stem_words = re.split(r"[_\-\s]+", stem)
        if fname.lower() in q_lower or stem in q_lower or (
            len(stem) > 3 and any(w in q_lower for w in stem_words if len(w) > 3)
        ):
            named_matches.append(fname)
    if named_matches:
        print(f"[RAG] Document resolution: named match -> {named_matches}")
        return named_matches

    # 2. Explicit multi-document intent -> everything uploaded
    if COMPARISON_PATTERN.search(query):
        print(f"[RAG] Document resolution: comparison intent -> all {filenames}")
        return filenames

    # 3. Default: files from the most recent upload action (batch of
    # several files if that's what was uploaded together, else just the
    # single most-recent file - unchanged from before for single uploads)
    if last_uploaded_batch:
        batch_valid = [f for f in last_uploaded_batch if f in uploaded_docs]
        if len(batch_valid) > 1:
            print(f"[RAG] Document resolution: defaulting to most recent batch -> {batch_valid}")
            return batch_valid
    if last_uploaded and last_uploaded in uploaded_docs:
        print(f"[RAG] Document resolution: defaulting to most recent -> {last_uploaded}")
        return [last_uploaded]

    print(f"[RAG] Document resolution: no last_uploaded tracked, falling back to all {filenames}")
    return filenames


def build_prompt(query, index, chat_history, memory_limit=6, extra_file_content="",
                  uploaded_docs=None, last_uploaded=None, last_uploaded_batch=None):
    if is_greeting_or_smalltalk(query):
        newline = "\n"
        recent_lines = newline.join(f"{m['role']}: {m['message']}" for m in chat_history[-4:])
        prompt = f"""Conversation so far:
{recent_lines}

User: {query}

Reply briefly and naturally (1-2 sentences)."""
        return prompt, [], False, True

    retrieved = retrieve(index, query)
    doc_context = "\n\n".join(chunk for chunk, src, score in retrieved)[:3000]
    uploaded_sources_used = []

    web_used = False
    web_context = ""
    if needs_web_search(query, retrieved):
        web_results = web_search(query)
        web_context = format_web_context(web_results)
        web_used = bool(web_context)

    if len(chat_history) > memory_limit:
        summary = summarize_messages(chat_history[:-memory_limit])
        recent_messages = chat_history[-memory_limit:]
        conversation_text = summary + "\n"
    else:
        recent_messages = chat_history
        conversation_text = ""

    for msg in recent_messages:
        conversation_text += f"{msg['role']}: {msg['message']}\n"

    conversation_text = conversation_text[-3000:]

    if uploaded_docs:
        # New multi-document path: resolve which uploaded file(s) this query
        # is actually about, then either hand over the FULL text (for
        # summary/overview requests, which need the whole document rather
        # than a handful of semantically-matched fragments) or the most
        # relevant chunks (for narrow factual questions), tagged per-source
        # so the model can cite which file each fact came from.
        resolved = resolve_referenced_documents(query, uploaded_docs, last_uploaded, last_uploaded_batch)
        uploaded_sources_used = resolved
        # Newest-first ordering: uploaded_docs is insertion-ordered (oldest
        # first), so when multiple docs are resolved (comparison mode) and
        # the total exceeds the character budget below, truncation trims
        # the OLDEST content first rather than silently squeezing out
        # whatever was just uploaded.
        resolved_ordered = sorted(
            resolved,
            key=lambda f: list(uploaded_docs.keys()).index(f),
            reverse=True
        )
        file_context_parts = []

        if SUMMARY_REQUEST_PATTERN.search(query):
            # Full-document mode. Cap per-doc and total length so we don't
            # blow the prompt budget if multiple large docs are resolved.
            per_doc_cap = 6000 if len(resolved_ordered) == 1 else 3000
            for fname in resolved_ordered:
                text = uploaded_docs.get(fname, "")
                file_context_parts.append(f"[Source: {fname}]\n{text[:per_doc_cap]}")
        else:
            # Targeted retrieval mode: chunk + embed each resolved doc and
            # pull the chunks most relevant to this specific question.
            q_vec = embed_texts([query])[0]
            for fname in resolved_ordered:
                text = uploaded_docs.get(fname, "")
                file_chunks = chunk_text(text)
                if not file_chunks:
                    continue
                file_embeddings = embed_texts(file_chunks)
                denom = np.linalg.norm(file_embeddings, axis=1) * np.linalg.norm(q_vec)
                denom[denom == 0] = 1e-8
                sims = (file_embeddings @ q_vec) / denom
                top_idx = np.argsort(sims)[::-1][:max(TOP_K, 5)]
                relevant_chunks = [file_chunks[i] for i in top_idx]
                chunk_block = "\n\n".join(relevant_chunks)[:3000]
                file_context_parts.append(f"[Source: {fname}]\n{chunk_block}")

        if file_context_parts:
            file_context = "\n\n".join(file_context_parts)[:8000]
            doc_context = (doc_context + "\n\n" + file_context).strip()

    elif extra_file_content:
        # Backward-compatible path for callers still passing a single
        # pre-joined string instead of the uploaded_docs dict.
        file_chunks = chunk_text(extra_file_content)
        if file_chunks:
            file_embeddings = embed_texts(file_chunks)
            q_vec = embed_texts([query])[0]
            denom = np.linalg.norm(file_embeddings, axis=1) * np.linalg.norm(q_vec)
            denom[denom == 0] = 1e-8
            sims = (file_embeddings @ q_vec) / denom
            top_idx = np.argsort(sims)[::-1][:max(TOP_K, 5)]
            relevant_chunks = [file_chunks[i] for i in top_idx]
            file_context = "\n\n".join(relevant_chunks)[:4000]
            doc_context = (doc_context + "\n\n" + file_context).strip()

    today_str = datetime.date.today().strftime("%A, %B %d, %Y")

    prompt = f"""Today's date is {today_str}.

Document Context:
{doc_context if doc_context else "(none)"}

Web Search Results:
{web_context if web_context else "(none)"}

Conversation so far:
{conversation_text}

User Question:
{query}

Give a complete, well-structured answer following your response style rules."""

    doc_sources = sorted(set(src for _, src, _ in retrieved) | set(uploaded_sources_used))
    return prompt, doc_sources, web_used, False


def chat_with_agent_stream(query, index, chat_history, memory_limit=6, extra_file_content="",
                            uploaded_docs=None, last_uploaded=None, last_uploaded_batch=None, model=None,
                            persona_prompt=None, memory_context=None):
    prompt, doc_sources, web_used, _ = build_prompt(
        query, index, chat_history, memory_limit, extra_file_content,
        uploaded_docs, last_uploaded, last_uploaded_batch
    )
    return stream_groq(prompt, model=model, persona_prompt=persona_prompt, memory_context=memory_context), doc_sources, web_used


RESEARCH_MAX_TOKENS = 3000


def build_research_prompt(query: str):
    """
    Runs two Tavily searches from slightly different angles (the raw query,
    plus a variant nudged toward analysis/comparison), de-dupes sources by
    URL, and builds a numbered source list the model is instructed to cite
    inline as [1], [2], etc. Returns (prompt, sources) or (None, []) if no
    web results came back at all (e.g. TAVILY_API_KEY missing/down).
    """
    today_str = datetime.date.today().strftime("%A, %B %d, %Y")
    r1 = web_search(query, max_results=6)
    r2 = web_search(f"{query} analysis comparison", max_results=6)

    sources = []
    seen_urls = set()
    for result in (r1, r2):
        if not result or result.get("_error"):
            continue
        for r in result.get("results", []):
            url = (r.get("url") or "").strip()
            if not url or url in seen_urls:
                continue
            seen_urls.add(url)
            sources.append({
                "title": (r.get("title") or url).strip(),
                "content": (r.get("content") or "").strip()[:600],
                "url": url,
            })
    sources = sources[:10]

    if not sources:
        return None, []

    source_block = "\n\n".join(
        f"[{i + 1}] {s['title']}\nURL: {s['url']}\nExcerpt: {s['content']}"
        for i, s in enumerate(sources)
    )

    prompt = f"""Today's date is {today_str}.

You are in AI RESEARCH MODE. Below are numbered sources from a live web search on this topic. Produce a structured research report:

1. Compare and cross-check claims across the sources - explicitly note where they agree and where they conflict.
2. Organize findings under clear headings.
3. Cite every factual claim inline using bracket numbers matching the source list, e.g. [1], [2].
4. Finish with a "References" section listing each numbered source's title and URL.
5. Structure: Overview, Key Findings, Conflicting Views (only if any), Conclusion, References.

Sources:
{source_block}

Research Question:
{query}

Write the full research report now."""
    return prompt, sources


def research_report_stream(query: str, model: str = None):
    prompt, sources = build_research_prompt(query)
    if prompt is None:
        def _no_sources():
            yield ("I couldn't pull web sources for this right now (search may be "
                   "unavailable or the topic returned nothing). Try rephrasing the "
                   "question or check back in a moment.")
        return _no_sources(), sources
    return stream_groq(prompt, max_tokens=RESEARCH_MAX_TOKENS, model=model), sources


# ============================================================================
# AI PRESENTATION STUDIO — free/open-source stack (Groq LLM + python-pptx).
# No paid APIs. Image fetching, Mermaid diagrams, and charts stay client-side
# in the frontend (already free: Mermaid.js / Chart.js / Pollinations.ai);
# this just turns model-authored JSON into a real, styled, downloadable
# .pptx file with editable text boxes, not a plain-text dump.
# ============================================================================

# NOTE: was 3500 - too small for larger decks (many slides / long bullets),
# so the model's JSON response was getting cut off mid-string by max_tokens,
# which made json.loads/raw_decode blow up with "Unterminated string
# starting at: ...". Raised the budget and also added a repair fallback
# below (_repair_truncated_json) so a cut-off response degrades gracefully
# instead of crashing the whole /generate-ppt request.
PPT_MAX_TOKENS = 6500

# A few built-in color themes so different templates don't all look identical.
PPT_THEMES = {
    "business": {
        "bg": (0x0F, 0x0F, 0x11), "accent": (0x7C, 0x5C, 0xFF),
        "text": (0xE8, 0xE8, 0xEA), "subtext": (0x9A, 0x9A, 0xA2),
    },
    "education": {
        "bg": (0x0C, 0x15, 0x12), "accent": (0x3D, 0xDC, 0x84),
        "text": (0xE6, 0xF2, 0xEC), "subtext": (0x8F, 0xB3, 0xA3),
    },
    "research": {
        "bg": (0x0B, 0x12, 0x20), "accent": (0x4F, 0x9D, 0xFF),
        "text": (0xE6, 0xED, 0xF7), "subtext": (0x8F, 0xA3, 0xC2),
    },
    "technical": {
        "bg": (0x28, 0x2A, 0x36), "accent": (0xBD, 0x93, 0xF9),
        "text": (0xF8, 0xF8, 0xF2), "subtext": (0xB3, 0xB6, 0xC4),
    },
}


def _rgb(t):
    return RGBColor(*t)


def build_slide_content_prompt(topic: str, n_slides: int, template: str) -> str:
    template_hint = {
        "business": "a business/pitch deck: problem, solution, market, product, "
                    "traction/roadmap, ask/next steps",
        "education": "an educational lecture deck: learning objectives, concept "
                     "explanations with examples, a recap/quiz slide",
        "research": "a research presentation: background, methodology, "
                    "findings/data, discussion, conclusion, references",
        "technical": "a technical/engineering deck: architecture overview, "
                     "components, data flow, tradeoffs, implementation notes",
    }.get(template, "a professional presentation")

    return f"""Return ONLY strict JSON, no markdown code fences, no commentary before or after.

Schema:
{{
  "title_slide": {{"title": "...", "subtitle": "..."}},
  "slides": [
    {{
      "type": "content" | "table" | "timeline" | "diagram",
      "title": "...",
      "bullets": ["...", "..."],
      "table": {{"headers": ["..."], "rows": [["...", "..."]]}},
      "timeline": [{{"label": "...", "detail": "..."}}],
      "diagram_mermaid": "flowchart TD; A-->B",
      "notes": "speaker notes for this slide, 1-3 sentences",
      "image_query": "2-4 word search term for a relevant royalty-free photo"
    }}
  ]
}}

Rules:
- Omit fields that don't apply to a slide's type (e.g. only "table" slides need "table").
- Every slide needs "title" and "notes". Content and diagram slides both need 6-8 detailed, informative bullets (roughly 15-25 words each, full sentences explaining the point, not short fragments) describing the same info the diagram would show.
- Include at least one "table" or "timeline" slide if the topic naturally fits one.
- Write exactly {n_slides} entries in "slides" (not counting the title slide).
- This is {template_hint}.

Topic: {topic}"""


def _repair_truncated_json(cleaned: str) -> dict:
    """
    Best-effort repair for when the LLM's JSON output got cut off mid-way
    (e.g. hit max_tokens mid-string, or the connection dropped). Closes any
    dangling open string, backs up to the last safe comma/bracket boundary,
    then closes any still-open braces/brackets in the right order. This
    means a truncated deck (say, 6 of 8 slides fully generated) still
    parses and renders instead of throwing json.decoder.JSONDecodeError
    and failing the whole /generate-ppt request.
    """
    s = cleaned

    # If we're inside an unterminated string (odd number of unescaped
    # quotes), drop back to the last quote so we don't leave a dangling
    # open quote confusing the bracket-balancing pass below.
    if s.count('"') % 2 == 1:
        s = s[:s.rfind('"') + 1]

    # Trim back to the last "safe" boundary - end of a complete key/value,
    # object, or array - so we don't try to keep a half-written field.
    last_safe = max(s.rfind(","), s.rfind("}"), s.rfind("]"))
    if last_safe != -1:
        s = s[:last_safe + 1] if s[last_safe] in "}]" else s[:last_safe]

    # Walk the string tracking bracket depth (ignoring brackets that appear
    # inside string literals) so we can close everything that's still open,
    # in the correct (reverse) order.
    opens = []
    in_str = False
    escaped = False
    for ch in s:
        if escaped:
            escaped = False
            continue
        if ch == "\\" and in_str:
            escaped = True
            continue
        if ch == '"':
            in_str = not in_str
            continue
        if in_str:
            continue
        if ch in "{[":
            opens.append(ch)
        elif ch in "}]":
            if opens:
                opens.pop()

    closers = {"{": "}", "[": "]"}
    s = s.rstrip().rstrip(",")
    s += "".join(closers[c] for c in reversed(opens))

    decoder = _json.JSONDecoder()
    obj, _ = decoder.raw_decode(s)
    return obj


def _parse_slide_json(raw: str) -> dict:
    cleaned = raw.strip()
    cleaned = re.sub(r"^```(json)?", "", cleaned.strip(), flags=re.IGNORECASE).strip()
    cleaned = re.sub(r"```$", "", cleaned.strip()).strip()
    start = cleaned.find("{")
    if start != -1:
        cleaned = cleaned[start:]
    decoder = _json.JSONDecoder()
    try:
        obj, _ = decoder.raw_decode(cleaned)
        return obj
    except _json.JSONDecodeError:
        # Response was cut off mid-JSON (usually hit max_tokens). Try to
        # salvage whatever complete slides we did get instead of failing
        # the whole request.
        return _repair_truncated_json(cleaned)


def generate_slide_content(topic: str, n_slides: int, template: str, model: str = None) -> dict:
    """Calls the LLM once and returns parsed slide-plan JSON."""
    prompt = build_slide_content_prompt(topic, n_slides, template)
    raw = "".join(stream_groq(prompt, max_tokens=PPT_MAX_TOKENS, model=model))
    data = None
    if raw.strip():
        try:
            data = _parse_slide_json(raw)
        except Exception:
            data = None
    if (not data or not data.get("slides")) and model != CHAT_MODEL:
        raw2 = "".join(stream_groq(prompt, max_tokens=PPT_MAX_TOKENS, model=CHAT_MODEL))
        try:
            data = _parse_slide_json(raw2)
        except Exception:
            pass
    if not data:
        data = _parse_slide_json(raw)
    return data


def _add_title_slide(prs, theme, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(theme["bg"])

    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(3.55), Inches(0.15), Inches(1.4))
    accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = _rgb(theme["accent"])
    accent_bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.6))
    tf = tbox.text_frame; tf.word_wrap = True
    tf.text = title or ""
    tf.paragraphs[0].font.size = Pt(40); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _rgb(theme["text"])

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.8))
        sf = sbox.text_frame; sf.word_wrap = True
        sf.text = subtitle
        sf.paragraphs[0].font.size = Pt(18)
        sf.paragraphs[0].font.color.rgb = _rgb(theme["subtext"])
    return slide


def _add_header(slide, theme, title):
    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.5), Inches(0.9))
    tf = tbox.text_frame
    tf.text = title or ""
    tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _rgb(theme["accent"])
    rule = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(3.2), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(theme["accent"])
    rule.line.fill.background()


def _add_content_slide(prs, theme, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _rgb(theme["bg"])
    _add_header(slide, theme, s.get("title", ""))

    body_width = Inches(8.2) if s.get("image_query") else Inches(11.5)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), body_width, Inches(5.3))
    bf = body.text_frame; bf.word_wrap = True
    bullets = s.get("bullets") or []
    for i, bullet in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = "•  " + str(bullet)
        p.font.size = Pt(20)
        p.font.color.rgb = _rgb(theme["text"])
        p.space_after = Pt(14)
    return slide


def _add_table_slide(prs, theme, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _rgb(theme["bg"])
    _add_header(slide, theme, s.get("title", ""))

    table_data = s.get("table") or {}
    headers = table_data.get("headers") or []
    rows = table_data.get("rows") or []
    if not headers:
        return _add_content_slide(prs, theme, s)

    n_rows, n_cols = len(rows) + 1, len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5 * n_rows)).table
    for c, h in enumerate(headers):
        cell = gtable.cell(0, c)
        cell.text = str(h)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(theme["accent"])
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(theme["bg"])
            cell.text_frame.paragraphs[0].font.color.rgb = _rgb(theme["text"])
    return slide


def _add_timeline_slide(prs, theme, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _rgb(theme["bg"])
    _add_header(slide, theme, s.get("title", ""))

    items = s.get("timeline") or []
    if not items:
        return _add_content_slide(prs, theme, s)

    n = len(items)
    total_w = 11.4
    step = total_w / n
    y_line = 3.3
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(y_line), Inches(total_w), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = _rgb(theme["accent"]); line.line.fill.background()

    for i, item in enumerate(items):
        x = 0.7 + i * step
        dot = slide.shapes.add_shape(9, Inches(x + step / 2 - 0.08), Inches(y_line - 0.06), Inches(0.16), Inches(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(theme["accent"]); dot.line.fill.background()

        lbl = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(step), Inches(0.7))
        lf = lbl.text_frame; lf.word_wrap = True
        lf.text = str(item.get("label", ""))
        lf.paragraphs[0].font.size = Pt(15); lf.paragraphs[0].font.bold = True
        lf.paragraphs[0].font.color.rgb = _rgb(theme["text"])

        det = slide.shapes.add_textbox(Inches(x), Inches(3.6), Inches(step), Inches(1.6))
        df = det.text_frame; df.word_wrap = True
        df.text = str(item.get("detail", ""))
        df.paragraphs[0].font.size = Pt(12)
        df.paragraphs[0].font.color.rgb = _rgb(theme["subtext"])
    return slide


def _add_illustration(slide, image_query: str, left, top, width, height):
    """Fetches a relevant image from Pollinations.ai and embeds it into slide. Silently no-ops on failure."""
    if not image_query:
        return
    try:
        url = f"https://image.pollinations.ai/prompt/{requests.utils.quote(image_query)}?width=800&height=800&nologo=true"
        resp = requests.get(url, timeout=25)
        resp.raise_for_status()
        slide.shapes.add_picture(io.BytesIO(resp.content), left, top, width=width, height=height)
    except Exception as e:
        print(f"[_add_illustration] Failed for '{image_query}': {e}")


def generate_presentation(topic: str, n_slides: int = 8, template: str = "business", model: str = None) -> bytes:
    """
    Full pipeline: LLM plans the deck -> python-pptx renders it into a real,
    editable, theme-styled .pptx, returned as raw bytes ready to stream back
    to the client. Diagram slides (diagram_mermaid) fall back to a bulleted
    content slide server-side, since Mermaid rendering happens client-side in
    the AI Studio panel / can be pasted into the notes for reference.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed on the server (pip install python-pptx).")

    theme = PPT_THEMES.get(template, PPT_THEMES["business"])
    data = generate_slide_content(topic, n_slides, template, model=model)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_info = data.get("title_slide") or {"title": topic, "subtitle": ""}
    _add_title_slide(prs, theme, title_info.get("title", topic), title_info.get("subtitle", ""))

    for s in data.get("slides", []):
        stype = (s.get("type") or "content").lower()
        if stype == "table":
            slide = _add_table_slide(prs, theme, s)
        elif stype == "timeline":
            slide = _add_timeline_slide(prs, theme, s)
        else:
            # "diagram" slides render as content slides server-side; the
            # Mermaid source (if any) is preserved in the speaker notes.
            slide = _add_content_slide(prs, theme, s)
            if stype == "diagram" and s.get("diagram_mermaid"):
                s["notes"] = (s.get("notes") or "") + "\n\nDiagram source:\n" + s["diagram_mermaid"]
            if s.get("image_query"):
                _add_illustration(slide, s["image_query"], Inches(9.3), Inches(1.6), Inches(3.4), Inches(3.4))

        notes = s.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def discover_topics():
    """
    Discover feed: returns 8-12 curated trending topics/questions daily.
    Stateless, no backend/DB needed - uses LLM to generate interesting
    topics on the fly. Frontend caches for 24 hours.
    """
    prompt = """Generate exactly 8 interesting, diverse trending topics/questions for an AI assistant user.
    Include mix of: tech news, DIY projects, learning topics, creative ideas, life hacks, productivity tips.
    Return as JSON array of objects: [{"title": "...", "desc": "..."}, ...]
    No preamble, just JSON."""
    
    try:
        raw = "".join(stream_groq(prompt, max_tokens=800))
        # Extract JSON from response
        start = raw.find('[')
        end = raw.rfind(']') + 1
        if start != -1 and end > start:
            topics = _json.loads(raw[start:end])
            return topics[:12] if isinstance(topics, list) else []
    except Exception as e:
        pass
    
    # Fallback hardcoded topics if LLM fails
    return [
        {"title": "Explain quantum computing simply", "desc": "Break down complex physics concepts"},
        {"title": "Python web scraping tutorial", "desc": "Build a data crawler from scratch"},
        {"title": "Best AI productivity tools 2026", "desc": "Discover the latest efficiency boosters"},
        {"title": "React hooks deep dive", "desc": "Master modern React patterns"},
        {"title": "Personal finance strategy", "desc": "Plan your financial future"},
        {"title": "Machine learning basics", "desc": "Start your ML journey today"},
        {"title": "Creative writing prompts", "desc": "Overcome writer's block instantly"},
        {"title": "System design interview prep", "desc": "Ace your next technical interview"},
    ]


def execute_plugin(plugin_name: str, query: str, url: str = None, method: str = "GET", headers_json: str = None):
    """
    Plugin execution: user-defined custom API endpoint.
    Frontend stores plugins as: {name, desc, url, method, headers_json}
    This function fetches, validates, calls the URL, and returns result.

    `url` is a template that may contain a {query} placeholder, which is
    replaced with the user's query text (URL-encoded) before the request
    is made.

    Returns: {"success": bool, "result": str, "error": str}
    """
    try:
        if not plugin_name or not query:
            return {"success": False, "result": "", "error": "Plugin name and query required"}

        if not url:
            return {
                "success": False,
                "result": "",
                "error": "No URL provided for this plugin."
            }

        import urllib.parse as _urlparse
        final_url = url.replace("{query}", _urlparse.quote(query))

        headers = {}
        if headers_json:
            try:
                headers = _json.loads(headers_json)
            except Exception:
                headers = {}

        method = (method or "GET").upper()
        if method == "GET":
            resp = requests.get(final_url, headers=headers, timeout=15)
        elif method == "POST":
            resp = requests.post(final_url, headers=headers, json={"query": query}, timeout=15)
        else:
            resp = requests.request(method, final_url, headers=headers, timeout=15)

        resp.raise_for_status()

        try:
            data = resp.json()
            result_text = _json.dumps(data)[:4000]
        except Exception:
            result_text = resp.text[:4000]

        return {"success": True, "result": result_text, "error": None}
    except requests.exceptions.RequestException as e:
        return {"success": False, "result": "", "error": f"Request failed: {e}"}
    except Exception as e:
        return {"success": False, "result": "", "error": str(e)}